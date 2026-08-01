"""
[checkba] PPTXBuilder 落字治理 + pptx_format_service 格式读写测试
"""
import os
import re
import zipfile

import pytest
from pptx import Presentation

from utils.pptx_builder import PPTXBuilder
from utils import pptx_format_utils as fmt
from utils.text_sanitizer import has_markdown
from services import pptx_format_service as pfs


@pytest.fixture
def built_pptx(tmp_path):
    """用 builder 生成一份含 markdown 输入的 pptx，返回路径"""
    b = PPTXBuilder()
    b.create_presentation()
    slide = b.add_blank_slide()
    b.add_text_element(slide, "# 一、公司治理结构", [50, 30, 700, 80], text_level=1)
    b.add_text_element(
        slide, "**核心结论**：本次交易*不构成*重大资产重组，~~初稿口径~~已废弃。",
        [50, 100, 700, 150])
    b.add_text_element(slide, "- 第一项要点\n- 第二项要点\n1. 编号项", [50, 160, 700, 300])
    b.add_table_element(
        slide,
        "<table><tr><th>股东</th><th>比例</th></tr><tr><td>**甲方**</td><td>51%</td></tr></table>",
        [50, 320, 700, 450])
    out = str(tmp_path / "built.pptx")
    b.save(out)
    return out


class TestBuilderMarkdownHygiene:
    def test_no_markdown_in_any_text(self, built_pptx):
        prs = Presentation(built_pptx)
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, 'has_text_frame', False):
                    assert not has_markdown(shape.text_frame.text), shape.text_frame.text
                if getattr(shape, 'has_table', False):
                    for row in shape.table.rows:
                        for cell in row.cells:
                            assert "**" not in cell.text

    def test_bold_converted_to_real_format(self, built_pptx):
        prs = Presentation(built_pptx)
        shape = list(prs.slides[0].shapes)[1]
        runs = shape.text_frame.paragraphs[0].runs
        bold_map = {r.text: r.font.bold for r in runs}
        assert bold_map.get("核心结论") is True
        italic_map = {r.text: r.font.italic for r in runs}
        assert italic_map.get("不构成") is True
        strike_runs = [r for r in runs if fmt.get_strike(r)]
        assert [r.text for r in strike_runs] == ["初稿口径"]

    def test_bullets_are_real_bullets(self, built_pptx):
        prs = Presentation(built_pptx)
        shape = list(prs.slides[0].shapes)[2]
        paras = shape.text_frame.paragraphs
        assert fmt.get_bullet(paras[0]) == "bullet"
        assert fmt.get_bullet(paras[1]) == "bullet"
        assert fmt.get_bullet(paras[2]) == "number"
        assert paras[0].text == "第一项要点"

    def test_house_fonts_written(self, built_pptx):
        with zipfile.ZipFile(built_pptx) as z:
            xml = z.read("ppt/slides/slide1.xml").decode("utf-8")
        assert f'<a:latin typeface="{fmt.HOUSE_LATIN_FONT}"' in xml
        assert f'<a:ea typeface="{fmt.HOUSE_EA_FONT}"' in xml

    def test_all_paragraphs_styled(self, built_pptx):
        """多行文本每一段都要有字号（原实现只有首段有）"""
        prs = Presentation(built_pptx)
        shape = list(prs.slides[0].shapes)[2]
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                assert run.font.size is not None


class TestFormatService:
    def test_inspect_structure(self, built_pptx):
        result = pfs.inspect(built_pptx)
        assert result["slide_count"] == 1
        shapes = result["slides"][0]["shapes"]
        assert any(s["has_table"] for s in shapes)
        text_shape = shapes[1]
        first_run = text_shape["paragraphs"][0]["runs"][0]
        assert first_run["font_name"] == fmt.HOUSE_LATIN_FONT
        assert first_run["ea_font"] == fmt.HOUSE_EA_FONT
        assert first_run["size_pt"] is not None

    def test_apply_run_and_paragraph_ops(self, built_pptx):
        ops = [
            {"action": "set_run_format", "slide": 0, "shape": 0,
             "format": {"strike": True, "highlight": "#FFFF00", "color": "#FF0000",
                        "size_pt": 30, "font_name": "Times New Roman", "ea_font": "黑体"}},
            {"action": "set_paragraph_format", "slide": 0, "shape": 0,
             "format": {"align": "center", "line_spacing": 1.5, "bullet": "none",
                        "space_after_pt": 6}},
        ]
        result = pfs.apply_ops(built_pptx, ops)
        assert result["failed"] == 0

        info = pfs.inspect(built_pptx)
        para = info["slides"][0]["shapes"][0]["paragraphs"][0]
        run = para["runs"][0]
        assert run["strike"] is True
        assert run["highlight"] == "#FFFF00"
        assert run["color"] == "#FF0000"
        assert run["size_pt"] == 30
        assert run["font_name"] == "Times New Roman"
        assert run["ea_font"] == "黑体"
        assert para["align"] == "center"
        assert para["line_spacing"] == 1.5
        assert para["space_after_pt"] == 6

    def test_replace_text_strips_markdown(self, built_pptx):
        ops = [{"action": "replace_text", "slide": 0, "shape": 1,
                "find": "核心结论", "replace": "**最终结论**"}]
        result = pfs.apply_ops(built_pptx, ops)
        assert result["failed"] == 0
        info = pfs.inspect(built_pptx)
        text = info["slides"][0]["shapes"][1]["paragraphs"][0]["text"]
        assert "最终结论" in text and "**" not in text

    def test_set_shape_text_converts_markdown(self, built_pptx):
        ops = [{"action": "set_shape_text", "slide": 0, "shape": 0,
                "text": "**重写标题**\n- 新要点甲\n- 新要点乙"}]
        result = pfs.apply_ops(built_pptx, ops)
        assert result["failed"] == 0
        prs = Presentation(built_pptx)
        tf = list(prs.slides[0].shapes)[0].text_frame
        assert tf.paragraphs[0].runs[0].text == "重写标题"
        assert tf.paragraphs[0].runs[0].font.bold is True
        assert fmt.get_bullet(tf.paragraphs[1]) == "bullet"
        assert tf.paragraphs[1].text == "新要点甲"

    def test_cell_ops(self, built_pptx):
        info = pfs.inspect(built_pptx)
        table_shape_idx = next(s["shape_index"] for s in info["slides"][0]["shapes"]
                               if s["has_table"])
        ops = [
            {"action": "set_cell_text", "slide": 0, "shape": table_shape_idx,
             "row": 1, "col": 0, "text": "**乙方**"},
            {"action": "set_cell_format", "slide": 0, "shape": table_shape_idx,
             "row": 1, "col": 0, "format": {"bold": True, "align": "left"}},
        ]
        result = pfs.apply_ops(built_pptx, ops)
        assert result["failed"] == 0
        prs = Presentation(built_pptx)
        table = next(s for s in prs.slides[0].shapes if getattr(s, 'has_table', False)).table
        cell = table.cell(1, 0)
        assert cell.text == "乙方"
        assert cell.text_frame.paragraphs[0].runs[0].font.bold is True

    def test_out_of_range_errors_reported(self, built_pptx):
        ops = [{"action": "set_run_format", "slide": 9, "shape": 0, "format": {"bold": True}}]
        result = pfs.apply_ops(built_pptx, ops)
        assert result["failed"] == 1
        assert "越界" in result["results"][0]["error"]

    def test_missing_file_raises(self):
        with pytest.raises(pfs.PptxFormatError):
            pfs.inspect("/nonexistent/x.pptx")


class TestApiEndpoints:
    def test_inspect_and_format_via_http(self, client, built_pptx):
        resp = client.post("/api/pptx/inspect", json={"pptx_path": built_pptx})
        assert resp.status_code == 200
        data = resp.get_json()["data"]
        assert data["slide_count"] == 1

        resp = client.post("/api/pptx/format", json={
            "pptx_path": built_pptx,
            "ops": [{"action": "set_run_format", "slide": 0, "shape": 0,
                     "format": {"bold": True}}],
        })
        assert resp.status_code == 200
        assert resp.get_json()["data"]["failed"] == 0

    def test_invalid_requests(self, client, built_pptx):
        assert client.post("/api/pptx/inspect", json={"pptx_path": "/no/x.pptx"}).status_code == 400
        assert client.post("/api/pptx/format", json={"pptx_path": built_pptx, "ops": []}).status_code == 400
