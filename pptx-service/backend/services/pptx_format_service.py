"""
[checkba] PPTX 格式识别与操作服务（python-pptx 层）。

给 AI 提供对存量 .pptx 文件与 docx 同级的格式能力：
- inspect: 全量结构化读出（形状/段落/run 的字体、字号、粗斜下划线、
  删除线、高亮、颜色、对齐、行距、项目符号；表格单元格）
- apply_ops: 批量格式操作（run 级/段落级/文本替换/整框重写/单元格写入），
  所有落字路径经过 text_sanitizer 去 markdown 化（标准格式落字）。

定位约定：slide/shape/paragraph/run/row/col 一律 0 起（与 inspect 输出一致）。
"""
import logging
import os
from typing import Any, Dict, List, Optional

from pptx import Presentation

from utils import pptx_format_utils as fmt
from utils.text_sanitizer import parse_lines, strip_markdown

logger = logging.getLogger(__name__)


class PptxFormatError(Exception):
    """格式操作错误（带上下文信息，控制器直接回给调用方）"""


def _load(pptx_path: str) -> Presentation:
    if not pptx_path or not os.path.exists(pptx_path):
        raise PptxFormatError(f"文件不存在: {pptx_path}")
    if not pptx_path.lower().endswith('.pptx'):
        raise PptxFormatError(f"仅支持 .pptx 文件: {pptx_path}")
    try:
        return Presentation(pptx_path)
    except Exception as e:
        raise PptxFormatError(f"无法解析 PPTX: {e}")


# ==================== inspect ====================

def inspect(pptx_path: str) -> Dict[str, Any]:
    prs = _load(pptx_path)
    slides = []
    for s_idx, slide in enumerate(prs.slides):
        shapes = []
        for sh_idx, shape in enumerate(slide.shapes):
            info: Dict[str, Any] = {
                'shape_index': sh_idx,
                'name': shape.name,
                'shape_type': str(shape.shape_type) if shape.shape_type is not None else None,
                'has_text_frame': bool(getattr(shape, 'has_text_frame', False)),
                'has_table': bool(getattr(shape, 'has_table', False)),
            }
            if info['has_text_frame']:
                info['paragraphs'] = _read_text_frame(shape.text_frame)
            if info['has_table']:
                tbl = shape.table
                info['table'] = {
                    'rows': len(tbl.rows),
                    'cols': len(tbl.columns),
                    'cells': [
                        {
                            'row': r, 'col': c,
                            'text': tbl.cell(r, c).text,
                            'paragraphs': _read_text_frame(tbl.cell(r, c).text_frame),
                        }
                        for r in range(len(tbl.rows))
                        for c in range(len(tbl.columns))
                    ],
                }
            shapes.append(info)
        slides.append({'slide_index': s_idx, 'shapes': shapes})
    return {'slide_count': len(prs.slides), 'slides': slides}


def _read_text_frame(tf) -> List[Dict[str, Any]]:
    paragraphs = []
    for p_idx, para in enumerate(tf.paragraphs):
        p_info = fmt.read_paragraph_format(para)
        p_info['paragraph_index'] = p_idx
        p_info['text'] = para.text
        p_info['runs'] = []
        for r_idx, run in enumerate(para.runs):
            r_info = fmt.read_run_format(run)
            r_info['run_index'] = r_idx
            p_info['runs'].append(r_info)
        paragraphs.append(p_info)
    return paragraphs


# ==================== apply_ops ====================

def apply_ops(pptx_path: str, ops: List[Dict[str, Any]],
              output_path: Optional[str] = None) -> Dict[str, Any]:
    """
    支持的 op（action 字段区分）：
    - set_run_format:      slide, shape, [paragraph], [run], format={run 级键}
                           缺省 run -> 该段全部 run；缺省 paragraph -> 全部段落
    - set_paragraph_format: slide, shape, [paragraph], format={段落级键}
    - replace_text:        slide, shape, find, replace（替换文本经 strip_markdown）
    - set_shape_text:      slide, shape, text（整框重写，markdown 转真格式，
                           可选 format={run 级键} 作为基础样式）
    - set_cell_text:       slide, shape, row, col, text（表格单元格，去 markdown）
    - set_cell_format:     slide, shape, row, col, [paragraph], format={run 级键 + 段落级键}
    返回每个 op 的执行结果。
    """
    prs = _load(pptx_path)
    results = []
    for i, op in enumerate(ops or []):
        action = op.get('action')
        try:
            handler = _OP_HANDLERS.get(action)
            if handler is None:
                raise PptxFormatError(f"未知 action: {action}")
            detail = handler(prs, op)
            results.append({'index': i, 'action': action, 'ok': True, 'detail': detail})
        except PptxFormatError as e:
            results.append({'index': i, 'action': action, 'ok': False, 'error': str(e)})
        except Exception as e:
            logger.exception(f"op #{i} ({action}) failed")
            results.append({'index': i, 'action': action, 'ok': False, 'error': str(e)})

    out = output_path or pptx_path
    prs.save(out)
    return {'output_path': out, 'applied': sum(1 for r in results if r['ok']),
            'failed': sum(1 for r in results if not r['ok']), 'results': results}


def _get_slide(prs, op):
    idx = op.get('slide')
    if idx is None or not (0 <= idx < len(prs.slides)):
        raise PptxFormatError(f"slide 越界: {idx}（共 {len(prs.slides)} 页）")
    return prs.slides[idx]


def _get_shape(prs, op):
    slide = _get_slide(prs, op)
    idx = op.get('shape')
    shapes = list(slide.shapes)
    if idx is None or not (0 <= idx < len(shapes)):
        raise PptxFormatError(f"shape 越界: {idx}（该页共 {len(shapes)} 个形状）")
    return shapes[idx]


def _get_text_frame(prs, op):
    shape = _get_shape(prs, op)
    if not getattr(shape, 'has_text_frame', False):
        raise PptxFormatError(f"shape {op.get('shape')} 没有文本框")
    return shape.text_frame


def _iter_paragraphs(tf, op):
    p_idx = op.get('paragraph')
    paras = list(tf.paragraphs)
    if p_idx is None:
        return paras
    if not (0 <= p_idx < len(paras)):
        raise PptxFormatError(f"paragraph 越界: {p_idx}（共 {len(paras)} 段）")
    return [paras[p_idx]]


def _op_set_run_format(prs, op):
    tf = _get_text_frame(prs, op)
    spec = op.get('format') or {}
    r_idx = op.get('run')
    count = 0
    for para in _iter_paragraphs(tf, op):
        runs = list(para.runs)
        targets = runs if r_idx is None else (
            [runs[r_idx]] if 0 <= r_idx < len(runs) else [])
        if r_idx is not None and not targets:
            raise PptxFormatError(f"run 越界: {r_idx}（该段共 {len(runs)} 个 run）")
        for run in targets:
            fmt.apply_run_format(run, spec)
            count += 1
    return {'runs_formatted': count}


def _op_set_paragraph_format(prs, op):
    tf = _get_text_frame(prs, op)
    spec = op.get('format') or {}
    count = 0
    for para in _iter_paragraphs(tf, op):
        fmt.apply_paragraph_format(para, spec)
        count += 1
    return {'paragraphs_formatted': count}


def _op_replace_text(prs, op):
    tf = _get_text_frame(prs, op)
    find = op.get('find')
    if not find:
        raise PptxFormatError("replace_text 需要非空 find")
    replace = strip_markdown(op.get('replace') or '')
    count = 0
    for para in tf.paragraphs:
        for run in para.runs:
            if find in run.text:
                run.text = run.text.replace(find, replace)
                count += 1
    if count == 0:
        raise PptxFormatError(f"未找到文本（run 级匹配）: {find}")
    return {'runs_replaced': count}


def _op_set_shape_text(prs, op):
    tf = _get_text_frame(prs, op)
    text = op.get('text')
    if text is None:
        raise PptxFormatError("set_shape_text 需要 text")
    base_spec = dict(op.get('format') or {})
    base_spec.setdefault('font_name', fmt.HOUSE_LATIN_FONT)
    base_spec.setdefault('ea_font', fmt.HOUSE_EA_FONT)

    tf.clear()
    line_specs = parse_lines(text)
    for idx, line in enumerate(line_specs):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p_fmt = {}
        if line.bullet:
            p_fmt['bullet'] = line.bullet
            if line.number:
                p_fmt['number_start'] = line.number
        if p_fmt:
            fmt.apply_paragraph_format(para, p_fmt)
        for ir in line.runs:
            run = para.add_run()
            run.text = ir.text
            spec = dict(base_spec)
            if ir.bold or line.heading_level:
                spec['bold'] = True
            if ir.italic:
                spec['italic'] = True
            if ir.strike:
                spec['strike'] = True
            if ir.highlight:
                spec['highlight'] = '#FFFF00'
            if ir.code:
                spec['code'] = True
            fmt.apply_run_format(run, spec)
    return {'paragraphs_written': len(line_specs)}


def _get_cell(prs, op):
    shape = _get_shape(prs, op)
    if not getattr(shape, 'has_table', False):
        raise PptxFormatError(f"shape {op.get('shape')} 不是表格")
    tbl = shape.table
    r, c = op.get('row'), op.get('col')
    if r is None or not (0 <= r < len(tbl.rows)):
        raise PptxFormatError(f"row 越界: {r}（共 {len(tbl.rows)} 行）")
    if c is None or not (0 <= c < len(tbl.columns)):
        raise PptxFormatError(f"col 越界: {c}（共 {len(tbl.columns)} 列）")
    return tbl.cell(r, c)


def _op_set_cell_text(prs, op):
    cell = _get_cell(prs, op)
    text = op.get('text')
    if text is None:
        raise PptxFormatError("set_cell_text 需要 text")
    cell.text = strip_markdown(text)
    for para in cell.text_frame.paragraphs:
        for run in para.runs:
            fmt.apply_run_format(run, {
                'font_name': fmt.HOUSE_LATIN_FONT,
                'ea_font': fmt.HOUSE_EA_FONT,
            })
    return {'text': cell.text}


def _op_set_cell_format(prs, op):
    cell = _get_cell(prs, op)
    spec = op.get('format') or {}
    run_keys = {'bold', 'italic', 'underline', 'strike', 'highlight',
                'color', 'font_name', 'ea_font', 'size_pt', 'code'}
    para_keys = {'align', 'line_spacing', 'space_before_pt', 'space_after_pt',
                 'bullet', 'number_start'}
    run_spec = {k: v for k, v in spec.items() if k in run_keys}
    para_spec = {k: v for k, v in spec.items() if k in para_keys}
    count = 0
    for para in _iter_paragraphs(cell.text_frame, op):
        if para_spec:
            fmt.apply_paragraph_format(para, para_spec)
        for run in para.runs:
            if run_spec:
                fmt.apply_run_format(run, run_spec)
                count += 1
    return {'runs_formatted': count}


_OP_HANDLERS = {
    'set_run_format': _op_set_run_format,
    'set_paragraph_format': _op_set_paragraph_format,
    'replace_text': _op_replace_text,
    'set_shape_text': _op_set_shape_text,
    'set_cell_text': _op_set_cell_text,
    'set_cell_format': _op_set_cell_format,
}
